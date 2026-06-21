#!/usr/bin/env python3
"""Build the midterm presentation deck (light academic style).

Embeds the real dashboard screenshots from docs/img/ and the verified
statistics from the data. Hand-built layout (no template) so it does not
read like a generic auto-generated slide pack.
"""
import pathlib
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = pathlib.Path(__file__).resolve().parents[1]
IMG = ROOT / "docs" / "img"
OUT = ROOT / "docs" / "midterm-presentation.pptx"

# ---- palette (warm off-white academic) ----
BG      = RGBColor(0xFC, 0xFB, 0xF8)
INK     = RGBColor(0x1B, 0x20, 0x27)
MUTE    = RGBColor(0x55, 0x5F, 0x69)
FAINT   = RGBColor(0x8A, 0x93, 0x9C)
TEAL    = RGBColor(0x0E, 0x9E, 0x8E)
TEAL_DK = RGBColor(0x0A, 0x6E, 0x62)
ORANGE  = RGBColor(0xCF, 0x55, 0x2E)
RULE    = RGBColor(0xDE, 0xE2, 0xDF)
PANEL   = RGBColor(0xF2, 0xF4, 0xF2)
TINT    = RGBColor(0xEA, 0xF3, 0xF1)
INK_PANEL = RGBColor(0x14, 0x16, 0x1A)  # for screenshot backing
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "PingFang SC"
MONO = "PingFang SC"

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Emu(int(SW * EMU_IN))
prs.slide_height = Emu(int(SH * EMU_IN))
BLANK = prs.slide_layouts[6]


def _set_font(run, size, bold=False, color=INK, name=FONT, italic=False, spacing=None):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)
    if spacing is not None:
        rPr.set("spc", str(spacing))


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, runs, first=False, space_after=4, space_before=0, line=1.0,
         align=PP_ALIGN.LEFT, bullet=False, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.line_spacing = line
    p.level = level
    if isinstance(runs, tuple):
        runs = [runs]
    for txt, kw in runs:
        r = p.add_run()
        r.text = txt
        _set_font(r, **kw)
    if bullet:
        _add_bullet(p)
    else:
        _no_bullet(p)
    return p


def _no_bullet(p):
    pPr = p._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum"):
        e = pPr.find(qn(tag))
        if e is not None:
            pPr.remove(e)
    if pPr.find(qn("a:buNone")) is None:
        pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def _add_bullet(p, char="—", color=TEAL):
    pPr = p._p.get_or_add_pPr()
    pPr.set("indent", "-182880")
    pPr.set("marL", "182880")
    buf = pPr.makeelement(qn("a:buFont"), {})
    buf.set("typeface", FONT)
    pPr.append(buf)
    bu = pPr.makeelement(qn("a:buChar"), {})
    bu.set("char", char)
    pPr.append(bu)


def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE,
         shadow=False, radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        _soft_shadow(sp)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def _soft_shadow(sp):
    spPr = sp._element.spPr
    el = spPr.makeelement(qn("a:effectLst"), {})
    sh = spPr.makeelement(qn("a:outerShdw"), {
        "blurRad": "90000", "dist": "38100", "dir": "5400000", "rotWithShape": "0"})
    clr = spPr.makeelement(qn("a:srgbClr"), {"val": "8A939C"})
    alpha = spPr.makeelement(qn("a:alpha"), {"val": "34000"})
    clr.append(alpha)
    sh.append(clr)
    el.append(sh)
    spPr.append(el)


def aspect(path):
    with Image.open(path) as im:
        return im.width / im.height


def picture(s, path, x, y, max_w, max_h, align="center", valign="top",
            frame=True, backing=False):
    """Fit image into box keeping aspect; optional thin frame + soft shadow."""
    ar = aspect(IMG / path if not isinstance(path, pathlib.Path) else path)
    w, h = max_w, max_w / ar
    if h > max_h:
        h, w = max_h, max_h * ar
    if align == "center":
        px = x + (max_w - w) / 2
    elif align == "right":
        px = x + (max_w - w)
    else:
        px = x
    if valign == "center":
        py = y + (max_h - h) / 2
    elif valign == "bottom":
        py = y + (max_h - h)
    else:
        py = y
    full = IMG / path if not isinstance(path, pathlib.Path) else path
    if frame:
        bg = rect(s, px - 0.04, py - 0.04, w + 0.08, h + 0.08, fill=INK_PANEL,
                  line=RULE, line_w=1.0, shadow=True)
    pic = s.shapes.add_picture(str(full), Inches(px), Inches(py), Inches(w), Inches(h))
    pic.line.color.rgb = RGBColor(0x2A, 0x2E, 0x34)
    pic.line.width = Pt(0.75)
    return px, py, w, h


def header(s, eyebrow, title, page, kicker_color=TEAL):
    # top hairline
    rect(s, 0.0, 0.0, SW, 0.10, fill=TEAL)
    _, tf = box(s, 0.62, 0.42, 11.6, 0.34)
    para(tf, (eyebrow.upper(), dict(size=11.5, bold=True, color=kicker_color,
         spacing=220)), first=True, space_after=0)
    _, tf = box(s, 0.62, 0.70, 12.1, 0.9)
    para(tf, (title, dict(size=26, bold=True, color=INK)), first=True, space_after=0)
    rect(s, 0.64, 1.46, 0.62, 0.045, fill=ORANGE)
    footer(s, page)


def footer(s, page):
    _, tf = box(s, 0.62, 7.08, 8.0, 0.3)
    para(tf, (f"LM Arena 榜单演化与机构格局可视分析  ·  中期汇报", dict(
        size=9, color=FAINT)), first=True, space_after=0)
    _, tf = box(s, 11.6, 7.08, 1.13, 0.3)
    para(tf, (f"{page:02d} / 15", dict(size=9, color=FAINT)), first=True,
         space_after=0, align=PP_ALIGN.RIGHT)


def callout(s, x, y, w, h, title, body, accent=TEAL, title_size=12.5, body_size=11):
    rect(s, x, y, w, h, fill=TINT if accent == TEAL else PANEL, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    rect(s, x, y, 0.075, h, fill=accent)
    _, tf = box(s, x + 0.26, y + 0.16, w - 0.42, h - 0.3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, (title, dict(size=title_size, bold=True, color=TEAL_DK if accent == TEAL
         else ORANGE)), first=True, space_after=4, line=1.05)
    if body:
        para(tf, (body, dict(size=body_size, color=INK)), space_after=0, line=1.18)


def chip(s, x, y, text, w=None, fill=PANEL, color=INK, bold=False, size=11):
    w = w or (0.26 + 0.105 * len(text))
    rect(s, x, y, w, 0.34, fill=fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.5)
    _, tf = box(s, x, y + 0.025, w, 0.30)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, (text, dict(size=size, bold=bold, color=color)), first=True,
         align=PP_ALIGN.CENTER, space_after=0)
    return w


# ============================================================
# Slide 1 — Title
# ============================================================
s = slide()
rect(s, 0.0, 0.0, SW, 0.14, fill=TEAL)
# left text column
_, tf = box(s, 0.85, 1.18, 6.5, 0.4)
para(tf, ("数据可视化 · 期中汇报", dict(size=13, bold=True, color=TEAL, spacing=180)),
     first=True, space_after=0)
_, tf = box(s, 0.83, 1.74, 6.6, 2.4)
para(tf, ("LM Arena 榜单演化", dict(size=40, bold=True, color=INK)), first=True,
     space_after=2, line=1.04)
para(tf, ("与机构格局可视分析", dict(size=40, bold=True, color=INK)), space_after=0,
     line=1.04)
rect(s, 0.86, 3.78, 0.9, 0.05, fill=ORANGE)
_, tf = box(s, 0.85, 3.98, 6.4, 1.2)
para(tf, ("从静态排行榜走向动态、可信、可解释的大模型比较", dict(size=15, color=MUTE)),
     first=True, space_after=10, line=1.25)
para(tf, [("数据集  ", dict(size=11.5, color=FAINT)),
          ("lmarena-ai/leaderboard-dataset", dict(size=11.5, bold=True, color=INK)),
          ("  ·  CC BY 4.0", dict(size=11.5, color=FAINT))], space_after=3)
para(tf, [("快照范围  ", dict(size=11.5, color=FAINT)),
          ("2023-05-08 → 2026-05-22", dict(size=11.5, bold=True, color=INK)),
          ("   ·   160 万行原始记录", dict(size=11.5, color=FAINT))], space_after=0)
_, tf = box(s, 0.85, 6.35, 6.4, 0.5)
para(tf, ("小组成员：____________      指导教师：____________", dict(size=11, color=FAINT)),
     first=True, space_after=0)
# right dashboard screenshot
picture(s, "00-hero-full.png", 7.35, 1.05, 5.5, 5.4, align="center", valign="center")
_, tf = box(s, 7.35, 6.42, 5.5, 0.3)
para(tf, ("实时仪表盘 · 单页多视图协调界面", dict(size=9.5, color=FAINT)), first=True,
     align=PP_ALIGN.CENTER, space_after=0)

# ============================================================
# Slide 2 — Design motivation
# ============================================================
s = slide()
header(s, "01 设计动机", "静态排行榜只回答“谁第一”，不回答“为什么可信”", 2)
_, tf = box(s, 0.62, 1.78, 7.0, 4.9)
para(tf, ("LMArena 是公开模型能力比较的重要入口，但原始榜单主要呈现某一时刻的名次。"
          "项目的设计动机不是把排行榜画得更漂亮，而是补上静态表格缺失的分析语境。",
          dict(size=13, color=INK)), first=True,
     space_after=12, line=1.3)
para(tf, [("但它并不是一张静态表，而是一个多维、随时间演化、带不确定性的数据立方：",
          dict(size=13, color=INK))], space_after=6, line=1.3)
para(tf, ("(arena × category × model × organization × snapshot)", dict(
     size=12.5, bold=True, color=TEAL_DK, name=MONO)), space_after=2, line=1.2)
para(tf, ("→ (rating, 置信区间, votes, rank)", dict(size=12.5, bold=True,
     color=TEAL_DK, name=MONO)), space_after=14, line=1.2)
para(tf, ("静态榜单会系统性地忽略四个关键问题：", dict(size=13, bold=True,
     color=INK)), space_after=8)
for t in ["头部领先权是否稳定，还是在频繁易主？",
          "相邻名次的分差，是否被置信区间淹没（不可信）？",
          "机构是少而精的深耕，还是多而平的铺量？",
          "综合榜的名次，是否掩盖了模型的跨类别偏科？"]:
    para(tf, (t, dict(size=12.5, color=INK)), space_after=6, line=1.2, bullet=True)
# right podium image + callout
picture(s, "01-podium.png", 7.95, 1.95, 4.78, 2.3, align="center", valign="top")
callout(s, 7.95, 4.55, 4.78, 1.95,
        "设计目标",
        "把这四个被静态榜单掩盖的维度变得可见、可探索、可解释——"
        "让用户不止记住“第一是谁”，而能回答“为什么可信、谁在追赶、适合什么任务”。",
        accent=ORANGE)

# ============================================================
# Slide 3 — Analysis tasks
# ============================================================
s = slide()
header(s, "02 任务定义", "把设计动机落到四类可操作分析任务", 3)
tasks = [
    ("T1  头部演化", "哪些模型在不同时间登顶？领先权是否稳定？", "排名演化 · 排名竞速"),
    ("T2  不确定性判读", "相邻名次的 rating 差异是否统计可信？", "榜单置信区间图"),
    ("T3  机构竞争", "机构是“少而精”还是“多而平”？", "机构位置图 · 覆盖图"),
    ("T4  模型偏科", "overall 是否掩盖了类别能力差异？", "详情抽屉 · 类别热力"),
]
ty = 1.82
for i, (name, q, view) in enumerate(tasks):
    yy = ty + i * 0.92
    rect(s, 0.62, yy, 7.05, 0.82, fill=PANEL, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    rect(s, 0.62, yy, 0.07, 0.82, fill=TEAL)
    _, tf = box(s, 0.92, yy + 0.10, 2.35, 0.62)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, (name, dict(size=13.5, bold=True, color=INK)), first=True, space_after=0)
    _, tf = box(s, 3.35, yy + 0.10, 3.05, 0.64)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, (q, dict(size=11, color=MUTE)), first=True, space_after=0, line=1.12)
    _, tf = box(s, 6.40, yy + 0.10, 1.20, 0.64)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, (view, dict(size=9.5, bold=True, color=TEAL_DK)), first=True,
         space_after=0, line=1.12, align=PP_ALIGN.RIGHT)
# right column: screenshot + why-vis callout
picture(s, "02-task-guide.png", 7.95, 1.86, 4.78, 1.55, align="center")
callout(s, 7.95, 3.46, 4.78, 1.02,
        "目标用户",
        "AI 应用开发者、模型能力研究者，以及需要展示可复现实验路径的课程答辩场景。",
        accent=ORANGE, body_size=10.5)
callout(s, 7.95, 4.62, 4.78, 1.95,
        "设计原则",
        "· 任务是探索式的——事先不知道“该问哪个问题”，单一统计量只能回答一个预设问题。\n"
        "· 数据带强不确定性（Bradley–Terry 区间），名次差常落在区间重叠内。\n"
        "· 关系是跨视图的——同一对象在机构、类别、时间轴上有多副面孔。",
        accent=TEAL, body_size=10.2)

# ============================================================
# Slide 4 — Data
# ============================================================
s = slide()
header(s, "03 数据说明", "一份多维、带不确定性的关系型数据立方", 4)
stats = [("1,601,272", "原始记录行"), ("20,521", "最新快照行"), ("14", "Arena"),
         ("571 / 559", "模型 (累计/最新)"), ("85", "机构"), ("46", "Category")]
sx, sw, gap = 0.62, 1.92, 0.05
for i, (num, lab) in enumerate(stats):
    xx = sx + i * (sw + gap)
    rect(s, xx, 1.82, sw, 1.0, fill=PANEL, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.09)
    _, tf = box(s, xx, 1.98, sw, 0.55)
    para(tf, (num, dict(size=19, bold=True, color=TEAL_DK)), first=True,
         align=PP_ALIGN.CENTER, space_after=0)
    _, tf = box(s, xx, 2.48, sw, 0.3)
    para(tf, (lab, dict(size=9.5, color=MUTE)), first=True, align=PP_ALIGN.CENTER,
         space_after=0)
# field table (left) + source note (right)
_, tf = box(s, 0.62, 3.18, 6.4, 0.4)
para(tf, ("核心字段", dict(size=13, bold=True, color=INK)), first=True, space_after=0)
fields = [
    ("model_name / organization", "模型与机构 — 搜索、联动高亮、下钻"),
    ("arena / category", "评测场景与类别 — 全局筛选与跨场景比较"),
    ("rating", "Bradley–Terry 评分 — 主要性能指标"),
    ("rating_lower / rating_upper", "置信区间上下界 — 不确定性可视化"),
    ("vote_count", "投票数 — 样本量与可靠性判断"),
    ("rank / leaderboard_publish_date", "名次与快照日期 — 时间演化与 brushing"),
]
fy = 3.62
for k, v in fields:
    rect(s, 0.62, fy, 0.09, 0.42, fill=TINT)
    _, tf = box(s, 0.86, fy + 0.015, 6.2, 0.46)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(k + "   ", dict(size=11, bold=True, color=INK, name=MONO)),
              (v, dict(size=10.5, color=MUTE))], first=True, space_after=0, line=1.1)
    fy += 0.50
callout(s, 7.35, 3.18, 5.38, 1.55,
        "数据来源",
        "Hugging Face · lmarena-ai/leaderboard-dataset（CC BY 4.0）。"
        "原始 parquet 可复现下载，不入库；前端只加载处理后的 compact JSON。",
        accent=TEAL, body_size=11)
callout(s, 7.35, 4.92, 5.38, 1.72,
        "口径说明",
        "latest split 表示各 arena/category 的最新可用记录；整体最晚日期为 2026-05-22，"
        "但不同 arena 的最新日期不完全相同。时间演化序列默认保留 Top-10。",
        accent=ORANGE, body_size=11)

# ============================================================
# Slide 5 — Data processing pipeline
# ============================================================
s = slide()
header(s, "03 数据处理流程", "160 万行原始数据 → 视图级 compact JSON", 5)
steps = ["下载 parquet\n(HuggingFace)", "标准化\n字段与日期",
         "计算 confidence\nwidth = upper−lower", "按视图\n预聚合",
         "7 个前端 JSON\n+ standalone 包"]
px0, pw, pgap, py0 = 0.62, 2.18, 0.34, 2.05
for i, st in enumerate(steps):
    xx = px0 + i * (pw + pgap)
    fill = TINT if i in (0, 4) else PANEL
    rect(s, xx, py0, pw, 1.18, fill=fill, line=RULE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    _, tf = box(s, xx + 0.12, py0, pw - 0.24, 1.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, (f"{i+1}", dict(size=12, bold=True, color=TEAL)), first=True,
         align=PP_ALIGN.CENTER, space_after=3)
    para(tf, (st, dict(size=11, bold=True, color=INK)), align=PP_ALIGN.CENTER,
         space_after=0, line=1.12)
    if i < len(steps) - 1:
        ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(xx + pw + 0.04),
                                Inches(py0 + 0.46), Inches(0.26), Inches(0.26))
        ar.fill.solid(); ar.fill.fore_color.rgb = TEAL
        ar.line.fill.background(); ar.shadow.inherit = False
# output JSON chips
_, tf = box(s, 0.62, 3.62, 12.1, 0.4)
para(tf, ("生成的前端数据产物", dict(size=12.5, bold=True, color=INK)), first=True,
     space_after=0)
outs = ["latest_leaderboard", "rank_timeseries_top", "organization_summary",
        "model_profiles", "date_category_counts", "arena_summary", "manifest"]
cx, cy = 0.62, 4.06
for o in outs:
    w = 0.26 + 0.105 * len(o)
    if cx + w > 12.6:
        cx, cy = 0.62, cy + 0.46
    chip(s, cx, cy, o, w=w, color=TEAL_DK, fill=PANEL, size=10.5)
    cx += w + 0.14
callout(s, 0.62, 5.05, 12.1, 1.45,
        "技术动机与结果",
        "难点：160 万行历史快照、28 个 parquet 文件、不同 arena 的最新日期不完全一致。"
        "策略：Python 在构建时完成清洗、补缺、confidence_width 派生和视图导向预聚合；"
        "结果：前端只读取 compact JSON，负责轻量筛选、联动与 SVG 绘制，并可打包为 standalone 离线演示。",
        accent=TEAL, body_size=11.5)

# ============================================================
# Slide 6 — System overview / innovation
# ============================================================
s = slide()
header(s, "04 核心创新 · 系统总览", "一套状态驱动的多视图探索系统，而不是几张孤立图", 6)
picture(s, "00-hero-full.png", 0.62, 1.78, 8.05, 5.0, align="left", valign="top")
_, tf = box(s, 8.95, 1.82, 3.78, 0.4)
para(tf, ("创新点如何落地", dict(size=13, bold=True, color=INK)), first=True, space_after=8)
regions = [("① 全局筛选", "Arena / Category / 机构 / Rank Top"),
           ("② Hero + 领奖台", "当前冠军、领先幅度、Top-3"),
           ("③ KPI 指标条", "随筛选/时间窗实时更新"),
           ("④ Bento 多视图", "7 个核心视图同屏协调"),
           ("⑤ 详情抽屉", "模型/机构 details-on-demand"),
           ("⑥ 联动状态栏", "时间窗、框选、搜索等状态可见可清除")]
for name, desc in regions:
    para(tf, [(name + "   ", dict(size=11.5, bold=True, color=TEAL_DK)),
              (desc, dict(size=10.5, color=MUTE))], space_after=7, line=1.12)
callout(s, 8.95, 5.55, 3.78, 1.18,
        "核心创新",
        "把“当前排名、历史变化、不确定性、机构格局、类别偏科”放进同一套交互状态中相互印证。",
        accent=ORANGE, body_size=10.5)

# ============================================================
# Slide 7 — Leaderboard / uncertainty (Case 2)
# ============================================================
s = slide()
header(s, "04 设计 · 当前榜单与置信区间", "把“名次是否可信”直接画进图形", 7)
picture(s, "04-leaderboard-uncertainty.png", 0.62, 1.78, 5.55, 5.0, align="left")
_, tf = box(s, 6.55, 1.82, 6.18, 1.7)
para(tf, ("视觉编码", dict(size=13, bold=True, color=INK)), first=True, space_after=7)
for t in [("x 位置 = rating（评分）", ),
          ("横线长度 = 置信区间（不确定性）", ),
          ("气泡大小 = vote_count（样本量）", ),
          ("颜色 = 机构", )]:
    para(tf, (t[0], dict(size=12, color=INK)), space_after=5, bullet=True, line=1.15)
callout(s, 6.55, 3.78, 6.18, 1.45,
        "实测发现（text / overall · 360 模型）",
        "置信区间宽度 与 log₁₀(票数) 的 Pearson 相关 = −0.91：票越多区间越窄。",
        accent=TEAL, body_size=11.5)
callout(s, 6.55, 5.38, 6.18, 1.35,
        "为什么 non-trivial",
        "低票模型区间宽达 ~42 Elo，高票收窄到 ~5 Elo；很多相邻名次落在区间重叠内，"
        "“第 8 vs 第 12”在统计上并不可区分。",
        accent=ORANGE, body_size=11)

# ============================================================
# Slide 8 — Evolution + race (Case 1)
# ============================================================
s = slide()
header(s, "04 设计 · 排名演化与竞速", "头部不是稳定霸权，而是高频更替的前沿竞赛", 8)
picture(s, "05-evolution.png", 0.62, 1.82, 4.35, 4.05, align="left", valign="top")
_, tf = box(s, 0.62, 5.95, 4.35, 0.4)
para(tf, ("排名演化 Bump chart · 可交互图例 / 十字准线", dict(size=9.5, color=FAINT)),
     first=True, align=PP_ALIGN.CENTER, space_after=0)
picture(s, "06-race.png", 5.15, 1.82, 7.55, 3.05, align="center", valign="top")
_, tf = box(s, 5.15, 4.92, 7.55, 0.35)
para(tf, ("排名竞速 Bar-chart race · 播放 / 拖拽 / 1×–4× 变速", dict(size=9.5,
     color=FAINT)), first=True, align=PP_ALIGN.CENTER, space_after=0)
callout(s, 5.15, 5.42, 7.55, 1.3,
        "实测发现（text / overall · 214 个历史快照）",
        "曾登顶 #1 的模型多达 22 个，经历 23 次易主，在 8 家机构间轮转"
        "（xAI → OpenAI → DeepSeek → Google → Anthropic …）。",
        accent=TEAL, body_size=11.5)

# ============================================================
# Slide 9 — Organization landscape (Case 4)
# ============================================================
s = slide()
header(s, "04 设计 · 机构格局", "“少而精” vs “多而平”，把广度与高度正交摆放", 9)
picture(s, "08-org-scatter.png", 0.62, 1.80, 4.85, 3.55, align="left", valign="top")
_, tf = box(s, 0.62, 5.42, 4.85, 0.3)
para(tf, ("机构位置图 · x=模型数, y=最好名次, 大小=总票数", dict(size=9, color=FAINT)),
     first=True, align=PP_ALIGN.CENTER, space_after=0)
picture(s, "09-org-bars.png", 5.62, 1.80, 4.05, 3.55, align="center", valign="top")
_, tf = box(s, 5.62, 5.42, 4.05, 0.3)
para(tf, ("机构模型覆盖（广度）", dict(size=9, color=FAINT)), first=True,
     align=PP_ALIGN.CENTER, space_after=0)
picture(s, "11-drawer-org.png", 9.85, 1.80, 2.88, 4.35, align="center", valign="top")
_, tf = box(s, 9.7, 6.30, 3.2, 0.3)
para(tf, ("机构详情抽屉 · 各 arena 表现", dict(size=9, color=FAINT)), first=True,
     align=PP_ALIGN.CENTER, space_after=0)
callout(s, 0.62, 5.80, 9.05, 0.95,
        "实测发现",
        "text/overall 前 20 名中 Google 5 + Anthropic 4 + OpenAI 4 = 13/20 由三家瓜分；"
        "散点图区分“研发深度”（左上少而精）与“机型铺量”（右下多而平）两种机构策略。",
        accent=TEAL, body_size=11)

# ============================================================
# Slide 10 — Model specialization drawer (Case 5)
# ============================================================
s = slide()
header(s, "04 设计 · 详情按需展开", "同一模型跨类别名次可差 100+ 位（偏科）", 10)
picture(s, "12-drawer-model.png", 0.62, 1.78, 3.0, 5.0, align="left", valign="top")
_, tf = box(s, 4.0, 1.86, 8.7, 1.9)
para(tf, ("details-on-demand：常驻右侧抽屉", dict(size=13, bold=True, color=INK)),
     first=True, space_after=7)
for t in ["点击任一模型 → 指标卡 + “类别强弱热力”逐条摊开",
          "每行一个 (arena / category)，色块按名次着色（青=强 · 红=弱）",
          "配合票数条，判断每个类别结论的样本量是否充分",
          "偏科模型呈青红混杂，通才模型整体同色——一眼可辨"]:
    para(tf, (t, dict(size=12, color=INK)), space_after=6, bullet=True, line=1.18)
callout(s, 4.0, 4.05, 8.7, 1.35,
        "实测发现（截图：mistral-small-3.1-24b）",
        "同一模型在不同 (arena/category) 名次从 #25 跨到 #74；综合榜的单一名次"
        "完全掩盖了这种能力结构差异。",
        accent=TEAL, body_size=11.5)
callout(s, 4.0, 5.55, 8.7, 1.2,
        "对“选型”的意义",
        "一个 overall 中游的模型，可能是某个垂直类别的强者——选型应看任务对口的类别，而非 overall。",
        accent=ORANGE, body_size=11.5)

# ============================================================
# Slide 11 — Core innovations
# ============================================================
s = slide()
header(s, "05 核心创新点", "从“看排名”升级为“解释排名、追踪变化、比较策略”", 11)
picture(s, "13-linking-leaderboard.png", 0.62, 1.80, 6.5, 3.45, align="left",
        valign="top")
_, tf = box(s, 0.62, 5.32, 6.5, 0.5)
para(tf, ("联动高亮：选中冠军后，全局同对象高亮、其余淡出（CSS class 切换，无重绘）",
     dict(size=9.5, color=FAINT)), first=True, align=PP_ALIGN.CENTER, space_after=0,
     line=1.1)
picture(s, "14-context-bar.png", 0.62, 5.95, 6.5, 0.6, align="left", valign="top")
_, tf = box(s, 7.4, 1.86, 5.3, 0.4)
para(tf, ("三层创新结构", dict(size=13, bold=True, color=INK)),
     first=True, space_after=8)
mech = [("不确定性进入主视图", "rating 不只画点，还同时编码置信区间与 vote_count，避免把不可靠名次当事实。"),
        ("时间演化可被操作", "排名演化 + bar-chart race + 时间窗 brushing，让“谁曾登顶、何时反超”可播放、可裁剪。"),
        ("机构策略二维化", "unique models × best rank 区分“少而精”和“多而平”，并支持框选多机构对比。"),
        ("全局状态协调", "7 个核心视图共享 state；点击、搜索、brush、详情抽屉互相印证，而非孤立刷新。"),
        ("按需细节解释", "模型/机构抽屉把 overall 背后的类别差异摊开，支撑偏科与选型判断。")]
for name, desc in mech:
    para(tf, [(name + " — ", dict(size=11.5, bold=True, color=TEAL_DK)),
              (desc, dict(size=10.5, color=MUTE))], space_after=8, line=1.12)

# ============================================================
# Slide 12 — Technical challenges
# ============================================================
s = slide()
header(s, "06 技术难点与解决方案", "性能、状态一致性、可复现性三件事同时成立", 12)
challenges = [
    ("C1  大规模历史数据无法直接进浏览器",
     "1,601,272 行 full 数据 + 28 个 parquet 文件；前端直接加载会慢、也缺 parquet 支持。",
     "Python 预处理：清洗字段、派生 confidence_width、按视图生成 compact JSON；时间序列默认 Top-10。"),
    ("C2  多视图交互容易状态冲突",
     "点击模型、框选机构、时间窗 brushing、搜索、全局筛选都可能同时存在。",
     "单一全局 state 管理 arena/category/org/focus/timeWindow/brush/search；联动状态栏显示可清除状态。"),
    ("C3  SVG 动画和高亮容易重绘/重影",
     "排名竞速播放时如果每帧重建 DOM，容易闪烁；半透明 dim 会出现拖影。",
     "竞速复用持久 row；高亮只切换 is-focus/is-dim class；时间窗变化才重建竞速数据。"),
    ("C4  离线答辩需要可复现演示",
     "现场网络与后端服务不稳定，原始数据也不应随 PPT 打包。",
     "standalone 构建把 HTML/CSS/JS/JSON 内联成单文件；README 与数据处理脚本可重新生成。"),
    ("边界  人类偏好原因还不能直接解释",
     "当前 leaderboard 只有聚合评分，没有每一次对战的 prompt/回答文本。",
     "本阶段先解释榜单结构与能力差异；两两胜负矩阵与对话案例面板作为后续扩展。"),
]
yy = 1.78
for i, (title, problem, solution) in enumerate(challenges):
    fill = TINT if i % 2 == 0 else PANEL
    rect(s, 0.62, yy, 12.1, 0.92, fill=fill, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    rect(s, 0.62, yy, 0.075, 0.92, fill=TEAL if i % 2 == 0 else ORANGE)
    _, tf = box(s, 0.88, yy + 0.10, 3.0, 0.72)
    para(tf, (title, dict(size=11.5, bold=True, color=INK)), first=True,
         space_after=0, line=1.12)
    _, tf = box(s, 3.95, yy + 0.10, 4.15, 0.72)
    para(tf, [("难点：", dict(size=9.7, bold=True, color=ORANGE)),
              (problem, dict(size=9.4, color=MUTE))], first=True,
         space_after=0, line=1.15)
    _, tf = box(s, 8.28, yy + 0.10, 4.2, 0.72)
    para(tf, [("解决：", dict(size=9.7, bold=True, color=TEAL_DK)),
              (solution, dict(size=9.4, color=INK))], first=True,
         space_after=0, line=1.15)
    yy += 1.02
callout(s, 0.62, 6.88, 12.1, 0.36,
        "实现状态",
        "已完成数据脚本、7 个核心视图、多视图联动、详情抽屉、排名竞速、时间窗 brushing、机构框选与 standalone 离线包。",
        accent=TEAL, title_size=9.6, body_size=9.2)

# ============================================================
# Slide 13 — Insights summary
# ============================================================
s = slide()
header(s, "07 初步洞察", "系统可复现的 5 个 non-trivial pattern", 13)
insights = [
    ("头部是高频更替的前沿竞赛", "214 个快照中 22 个模型曾登顶、23 次易主、跨 8 家机构", "排名演化 + 竞速"),
    ("评分不确定性几乎由票数决定", "区间宽度 vs log₁₀(票数) Pearson = −0.91，相邻名次常不可区分", "榜单置信区间"),
    ("风格控制能把名次搬动 50–70 位", "去除文风红利后，部分模型大涨、部分大跌，暴露排版偏好", "arena 切换 + 搜索"),
    ("机构格局是寡头 + 长尾", "前 20 名 13/20 归三家；少而精 vs 多而平两类策略分化", "机构位置图 + 抽屉"),
    ("存在显著偏科模型", "同一模型跨类别名次可差 100+ 位，overall 掩盖能力结构", "详情类别热力"),
]
iy = 1.84
for i, (title, body, view) in enumerate(insights):
    rect(s, 0.62, iy, 12.1, 0.92, fill=PANEL if i % 2 else TINT, line=None,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    _, tf = box(s, 0.86, iy + 0.085, 0.6, 0.74)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, (f"{i+1}", dict(size=22, bold=True, color=TEAL)), first=True,
         space_after=0)
    _, tf = box(s, 1.55, iy + 0.10, 7.6, 0.74)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, (title, dict(size=13, bold=True, color=INK)), first=True, space_after=2,
         line=1.05)
    para(tf, (body, dict(size=10.5, color=MUTE)), space_after=0, line=1.1)
    _, tf = box(s, 9.35, iy + 0.10, 3.15, 0.74)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, (view, dict(size=10, bold=True, color=TEAL_DK)), first=True,
         align=PP_ALIGN.RIGHT, space_after=0, line=1.1)
    iy += 1.0

# ============================================================
# Slide 14 — Remaining work
# ============================================================
s = slide()
header(s, "08 待完成工作", "从“可用”走向“好讲、可辩、可复现”", 14)
left = [("Demo 与演示流程", "固化现场操作脚本，确保每个发现都能按路径复现。"),
        ("Case study 打磨", "为 5 个发现配标注截图，明确观察顺序与结论边界。"),
        ("数据质量提示", "补充缺失值、小样本与 latest 日期差异说明，避免误读。")]
right = [("授权类型口径", "license 字段可做初步开源/闭源分析，但严格判断需额外规则。"),
         ("偏好原因扩展", "如时间允许，引入 human preference 数据做两两胜负矩阵与对话案例面板。"),
         ("复现性终检", "核对 README、数据处理脚本、standalone 包与 AI 使用声明。")]
for col, items in ((0.62, left), (6.75, right)):
    yy = 1.95
    for title, body in items:
        rect(s, col, yy, 5.95, 1.45, fill=PANEL, line=None,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
        rect(s, col, yy, 0.07, 1.45, fill=ORANGE if col > 1 else TEAL)
        _, tf = box(s, col + 0.28, yy + 0.18, 5.5, 1.12)
        para(tf, (title, dict(size=13, bold=True, color=INK)), first=True,
             space_after=5)
        para(tf, (body, dict(size=11, color=MUTE)), space_after=0, line=1.22)
        yy += 1.62

# ============================================================
# Slide 15 — Defense points / closing
# ============================================================
s = slide()
rect(s, 0.0, 0.0, SW, 0.14, fill=TEAL)
_, tf = box(s, 0.85, 1.05, 11.6, 0.4)
para(tf, ("答辩要点", dict(size=13, bold=True, color=TEAL, spacing=200)),
     first=True, space_after=0)
_, tf = box(s, 0.83, 1.5, 11.6, 0.8)
para(tf, ("围绕设计动机、核心创新与技术难点的答辩主线", dict(
     size=25, bold=True, color=INK)), first=True, space_after=0)
rect(s, 0.86, 2.42, 0.8, 0.05, fill=ORANGE)
points = [
    "设计动机：静态榜单只给出名次，本系统补上可信度、时间演化、机构策略与类别偏科。",
    "核心创新：7 个核心视图共享同一状态，让排名、区间、时间、机构、详情互相解释。",
    "技术难点：160 万行历史数据、复杂交互状态、SVG 动画性能，通过预聚合与状态机解决。",
    "可复现性：数据处理脚本、public/data JSON、standalone HTML 构成完整离线演示链路。",
    "结论价值：5 个 non-trivial pattern 都能通过明确交互路径当场复现，而不是只停留在截图。",
]
_, tf = box(s, 0.85, 2.85, 11.7, 3.6)
for i, p in enumerate(points):
    para(tf, [(f"{i+1}   ", dict(size=14, bold=True, color=TEAL)),
              (p, dict(size=13.5, color=INK))], first=(i == 0), space_after=13,
         line=1.25)
_, tf = box(s, 0.85, 6.6, 11.7, 0.4)
para(tf, ("谢谢 · 欢迎提问", dict(size=12, color=FAINT)), first=True, space_after=0)
footer(s, 15)

prs.save(str(OUT))
print("saved", OUT, "·", len(prs.slides._sldIdLst), "slides")
